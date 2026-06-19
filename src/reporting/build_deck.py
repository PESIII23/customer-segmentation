"""
Assemble the stakeholder business presentation (.pptx).

Embeds the business-friendly charts from docs/deck_assets/ and selected EDA
artifacts from docs/, and writes the deck to docs/Quote_Decision_Predictor_Business_Presentation.pptx.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "docs"
ASSETS = DOCS / "deck_assets"
OUT = DOCS / "Quote_Decision_Predictor_Business_Presentation.pptx"

# Palette
NAVY = RGBColor(0x0B, 0x25, 0x45)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
AMBER = RGBColor(0xE9, 0xC4, 0x6A)
RUST = RGBColor(0xE7, 0x6F, 0x51)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x5A, 0x6B, 0x7B)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _solid(sh, color)
    sh.shadow.inherit = False
    return sh


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {text, size, color, bold, space_after, bullet}."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        if ln.get("space_before") is not None:
            p.space_before = Pt(ln["space_before"])
        run = p.add_run()
        run.text = ("•  " if ln.get("bullet") else "") + ln["text"]
        f = run.font
        f.size = Pt(ln.get("size", 18))
        f.bold = ln.get("bold", False)
        f.color.rgb = ln.get("color", NAVY)
        f.name = "Calibri"
    return tb


def header(slide, kicker, title):
    rect(slide, 0, 0, SW, Inches(1.25), NAVY)
    rect(slide, 0, Inches(1.25), SW, Inches(0.06), TEAL)
    textbox(slide, Inches(0.6), Inches(0.12), Inches(12), Inches(0.4),
            [{"text": kicker, "size": 13, "color": AMBER, "bold": True}])
    textbox(slide, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.75),
            [{"text": title, "size": 27, "color": WHITE, "bold": True}])


def footer(slide, n):
    textbox(slide, Inches(0.6), Inches(7.08), Inches(8), Inches(0.35),
            [{"text": "Quote Decision Predictor  ·  Confidential — Internal",
              "size": 10, "color": GREY}])
    textbox(slide, Inches(12.2), Inches(7.08), Inches(0.8), Inches(0.35),
            [{"text": str(n), "size": 10, "color": GREY, "align": PP_ALIGN.RIGHT}])


def pic_fit(slide, path, x, y, w, h):
    """Add picture scaled to fit within (w,h), centered in the box."""
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
    except Exception:
        slide.shapes.add_picture(str(path), x, y, width=w)
        return
    ratio = min(w / iw, h / ih)
    nw, nh = int(iw * ratio), int(ih * ratio)
    px = x + int((w - nw) / 2)
    py = y + int((h - nh) / 2)
    slide.shapes.add_picture(str(path), px, py, width=Emu(nw), height=Emu(nh))


def card(slide, x, y, w, h, head, body, accent=TEAL):
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, Inches(0.09), h, accent)
    lines = [{"text": head, "size": 16, "bold": True, "color": NAVY,
              "space_after": 6}]
    for b in body:
        lines.append({"text": b, "size": 13, "color": GREY, "space_after": 4})
    textbox(slide, x + Inches(0.28), y + Inches(0.16), w - Inches(0.45),
            h - Inches(0.3), lines)


# ---------------------------------------------------------------- 1. Title
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.62), SW, Inches(0.08), TEAL)
textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
        [{"text": "GEICO  ·  AUTO QUOTE ANALYTICS", "size": 16,
          "color": AMBER, "bold": True}])
textbox(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(1.8),
        [{"text": "Predicting Quote Decisions", "size": 46, "color": WHITE,
          "bold": True, "space_after": 6},
         {"text": "Turning 5,000 historic quotes into a conversion playbook",
          "size": 22, "color": RGBColor(0xCF, 0xDD, 0xE8)}])
textbox(s, Inches(0.9), Inches(4.85), Inches(11), Inches(1.4),
        [{"text": "Business stakeholder presentation — Sales & Underwriting",
          "size": 16, "color": WHITE, "space_after": 4},
         {"text": "Prepared by Phillip Smith  ·  Senior ML Engineering",
          "size": 14, "color": GREY, "space_after": 4},
         {"text": "June 2026", "size": 14, "color": GREY}])

# ---------------------------------------------------------- 2. Use case / objective
s = prs.slides.add_slide(BLANK)
header(s, "THE BUSINESS QUESTION", "Use Case & Objective")
textbox(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.9),
        [{"text": "When a prospect requests an auto quote, can we tell up front "
          "how likely they are to accept — and what drives that decision?",
          "size": 18, "color": NAVY, "bold": True}])
card(s, Inches(0.6), Inches(2.65), Inches(3.9), Inches(2.0), "Objective",
     ["Predict whether a quote ends in Accepted or Denied",
      "Identify the levers that move conversion",
      "Give Sales a data-backed targeting & follow-up guide"], TEAL)
card(s, Inches(4.7), Inches(2.65), Inches(3.9), Inches(2.0), "Why it matters",
     ["68.6% of quotes convert today",
      "Even a few points of lift = meaningful premium",
      "Focus effort where it actually pays off"], AMBER)
card(s, Inches(8.8), Inches(2.65), Inches(3.9), Inches(2.0), "Success criteria",
     ["Beat the 'accept everyone' baseline",
      "Errors acceptable for a sales-prioritization use",
      "Explainable to underwriting & sales"], RUST)
card(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(1.5),
     "Scope & framing",
     ["This is a CLASSIFICATION problem (Accept vs. Deny) — not a pricing model. "
      "The data contains the decision outcome but no premium/price field, so margin "
      "impact is discussed qualitatively. Built on 5,000 historic quotes with a "
      "held-out 20% test set so results reflect unseen prospects."], NAVY)
footer(s, 2)

# ---------------------------------------------------------- 3. Data snapshot
s = prs.slides.add_slide(BLANK)
header(s, "WHAT WE'RE WORKING WITH", "Data Snapshot")
stats = [("5,000", "historic quotes", TEAL), ("15", "applicant & vehicle attributes", AMBER),
         ("68.6%", "overall acceptance rate", RUST), ("0", "missing values", NAVY)]
x = Inches(0.6)
for val, lab, col in stats:
    rect(s, x, Inches(1.55), Inches(2.9), Inches(1.5), LIGHT)
    rect(s, x, Inches(1.55), Inches(2.9), Inches(0.12), col)
    textbox(s, x, Inches(1.78), Inches(2.9), Inches(0.8),
            [{"text": val, "size": 38, "color": col, "bold": True,
              "align": PP_ALIGN.CENTER}])
    textbox(s, x + Inches(0.15), Inches(2.55), Inches(2.6), Inches(0.45),
            [{"text": lab, "size": 12, "color": GREY, "align": PP_ALIGN.CENTER}])
    x += Inches(3.05)
card(s, Inches(0.6), Inches(3.35), Inches(5.95), Inches(3.2),
     "What each quote tells us",
     ["Applicant: age, gender, marital & home status, education",
      "Behavior: driving experience, risk tier, daily / annual mileage",
      "Vehicle: car type, car age, ownership, number of vehicles",
      "Quote: type (Auto-only vs. bundled) and the final decision",
      "",
      "Target = QuoteSts  →  Accepted (1) or Denied (0)"], TEAL)
rect(s, Inches(6.75), Inches(3.35), Inches(5.95), Inches(3.2), LIGHT)
textbox(s, Inches(7.0), Inches(3.5), Inches(5.5), Inches(0.5),
        [{"text": "Outcome mix — a healthy, slightly imbalanced base",
          "size": 15, "bold": True, "color": NAVY}])
pic_fit(s, DOCS / "eda_countplots" / "STAT_Q_cnt.png",
        Inches(6.95), Inches(3.95), Inches(5.6), Inches(2.45))
footer(s, 3)

# ---------------------------------------------------------- 4. EDA headline drivers
s = prs.slides.add_slide(BLANK)
header(s, "EDA FINDINGS — IN PLAIN TERMS", "Two Factors Dominate the Decision")
pic_fit(s, ASSETS / "rate_quotetype.png", Inches(0.5), Inches(1.5),
        Inches(6.1), Inches(3.3))
pic_fit(s, ASSETS / "rate_sex.png", Inches(6.7), Inches(1.5),
        Inches(6.1), Inches(3.3))
card(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(1.5),
     "So what?",
     ["Bundled (Auto + another policy) quotes convert at 80% vs. 52% for auto-only — "
      "a 28-point gap. Male applicants accept at 81% vs. 58% for female applicants. "
      "These two signals alone explain most of the conversion story."], RUST)
footer(s, 4)

# ---------------------------------------------------------- 5. EDA secondary + age
s = prs.slides.add_slide(BLANK)
header(s, "EDA FINDINGS — SECONDARY SIGNALS", "Most Other Factors Barely Move the Needle")
pic_fit(s, ASSETS / "rate_driverisk.png", Inches(0.5), Inches(1.5),
        Inches(6.1), Inches(3.2))
pic_fit(s, ASSETS / "rate_cartype.png", Inches(6.7), Inches(1.5),
        Inches(6.1), Inches(3.2))
card(s, Inches(0.6), Inches(4.85), Inches(5.95), Inches(1.6),
     "The long tail of weak factors",
     ["Driving-risk tier, car type, mileage, education, marital and home status "
      "all hover within a few points of the 69% average — interesting, but weak "
      "predictors on their own."], AMBER)
card(s, Inches(6.75), Inches(4.85), Inches(5.95), Inches(1.6),
     "Age: volume, not discrimination",
     ["74% of applicants are 17–25. Acceptance is fairly flat across age bands, "
      "so age tells us where the customers are more than who converts."], TEAL)
footer(s, 5)

# ---------------------------------------------------------- 6. Data prep
s = prs.slides.add_slide(BLANK)
header(s, "GETTING THE DATA MODEL-READY", "Data Preparation & Wrangling")
cards = [
    ("Standardize & clean", ["Renamed 17 raw columns to clear names",
                             "Confirmed zero missing values",
                             "Raw file kept immutable; cached to Parquet"], TEAL),
    ("Encode categories", ["Yes/No fields → 0/1 (gender, quote type, ownership)",
                           "Ordered tiers → ranked numbers (risk, mileage, education)",
                           "Unordered types → one-hot (car type, marital status)"], AMBER),
    ("Binning", ["Grouped raw ages into bands (17–25, 26–35, …)",
                 "Made age patterns readable for stakeholders",
                 "Used for EDA storytelling, not as a model crutch"], RUST),
    ("Feature selection", ["Chi-square test ranks each factor vs. outcome",
                           "Kept the 6 statistically significant signals",
                           "Dropped 12 fields that added only noise"], NAVY),
]
x = Inches(0.6)
for head_t, body, col in cards:
    card(s, x, Inches(1.6), Inches(2.92), Inches(3.6), head_t, body, col)
    x += Inches(3.04)
card(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.15),
     "Leakage guardrail",
     ["Scaling is fit inside the training fold only — never on the full dataset — "
      "so test results honestly reflect performance on brand-new quotes."], TEAL)
footer(s, 6)

# ---------------------------------------------------------- 7. Approach
s = prs.slides.add_slide(BLANK)
header(s, "HOW WE TESTED THE QUESTION", "Analytical Approach")
steps = [("1", "Traditional analytics", "Acceptance rates, cross-tabs, and a simple "
          "2-variable business rule (bundled OR male → predict accept)"),
         ("2", "Machine learning", "Three classifiers — Logistic Regression, "
          "Random Forest, Gradient Boosting — on the 6 selected features"),
         ("3", "Fair evaluation", "80/20 stratified split; judged on accuracy, "
          "precision, recall and ROC-AUC against unseen quotes"),
         ("4", "Head-to-head", "Compared ML to the simple rule and to the "
          "'accept everyone' baseline to isolate true added value")]
y = Inches(1.65)
for num, t, d in steps:
    rect(s, Inches(0.6), y, Inches(0.85), Inches(1.05), NAVY)
    textbox(s, Inches(0.6), y + Inches(0.22), Inches(0.85), Inches(0.6),
            [{"text": num, "size": 30, "color": AMBER, "bold": True,
              "align": PP_ALIGN.CENTER}])
    rect(s, Inches(1.6), y, Inches(11.1), Inches(1.05), LIGHT)
    textbox(s, Inches(1.85), y + Inches(0.13), Inches(10.6), Inches(0.85),
            [{"text": t, "size": 17, "bold": True, "color": NAVY,
              "space_after": 2},
             {"text": d, "size": 13, "color": GREY}])
    y += Inches(1.2)
footer(s, 7)

# ---------------------------------------------------------- 8. Model results
s = prs.slides.add_slide(BLANK)
header(s, "MODEL OUTPUT — THE HEADLINE", "Results: Good, But Not Because of ML")
pic_fit(s, ASSETS / "model_comparison.png", Inches(0.5), Inches(1.5),
        Inches(7.6), Inches(4.0))
card(s, Inches(8.35), Inches(1.6), Inches(4.4), Inches(1.85),
     "Best model: Gradient Boosting",
     ["Accuracy 78.5%   ·   ROC-AUC 0.72",
      "Recall 91%  (catches most true accepts)",
      "Precision 80%  ·  F1 0.85"], NAVY)
card(s, Inches(8.35), Inches(3.6), Inches(4.4), Inches(1.95),
     "The honest finding",
     ["A 2-variable rule hits 78.3% — within 0.2 pts of the best ML model. "
      "Almost all the signal lives in two factors any analyst can read."], RUST)
textbox(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.1),
        [{"text": "Verified on held-out test data (1,000 quotes the models never saw). "
          "Lift over 'accept everyone' is ~10 points — but ML's lift over simple "
          "analytics is essentially zero.", "size": 13, "color": GREY}])
footer(s, 8)

# ---------------------------------------------------------- 9. Errors
s = prs.slides.add_slide(BLANK)
header(s, "ARE THE MISTAKES ACCEPTABLE?", "Type I & Type II Errors")
pic_fit(s, ASSETS / "confusion_gb.png", Inches(0.5), Inches(1.55),
        Inches(6.0), Inches(4.6))
card(s, Inches(6.8), Inches(1.65), Inches(5.9), Inches(1.55),
     "Type II — missed accepts (61)",
     ["Model says 'deny' but they'd have accepted. The bigger business cost: "
      "under-served warm prospects. Kept low — recall is 91%."], RUST)
card(s, Inches(6.8), Inches(3.35), Inches(5.9), Inches(1.55),
     "Type I — false alarms (154)",
     ["Model says 'accept' but they deny. Low cost here: a follow-up call to "
      "someone who won't convert. Tolerable for sales prioritization."], AMBER)
card(s, Inches(6.8), Inches(5.05), Inches(5.9), Inches(1.4),
     "Verdict",
     ["For prioritizing sales effort, this error balance is acceptable. We "
      "deliberately favor recall — better to over-include warm leads than miss them."], TEAL)
footer(s, 9)

# ---------------------------------------------------------- 10. ML vs analytics reflection
s = prs.slides.add_slide(BLANK)
header(s, "REFLECTION", "Did Machine Learning Add Value Here?")
card(s, Inches(0.6), Inches(1.6), Inches(5.95), Inches(2.35),
     "Incremental value: minimal",
     ["ML beat 'accept everyone' by ~10 pts — but matched the simple rule",
      "The dataset's signal is concentrated in 2 obvious factors",
      "No complex interactions for trees to exploit",
      "ROC-AUC 0.72 = useful ranking, not a step-change"], RUST)
card(s, Inches(6.7), Inches(1.6), Inches(6.0), Inches(2.35),
     "Cost / benefit",
     ["Simple analytics: hours, fully transparent, easy to govern",
      "ML pipeline: more build & maintenance, harder to explain",
      "Here the ROI of ML over analytics is hard to justify",
      "ML earns its keep when data is richer & relationships non-linear"], NAVY)
card(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.25),
     "What I learned about analytics vs. machine learning",
     ["Analytics explains the 'why' and is faster to trust; ML scales pattern-finding "
      "when relationships are too complex to eyeball.",
      "More model complexity ≠ more value — match the tool to the signal in the data.",
      "Did binning / prep help? Prep was essential for clean modeling and clear EDA, "
      "but binning aided human storytelling more than model accuracy.",
      "Always benchmark ML against a dumb baseline AND a simple rule — or you can't "
      "tell whether the model is actually earning its complexity."], TEAL)
footer(s, 10)

# ---------------------------------------------------------- 11. Business value
s = prs.slides.add_slide(BLANK)
header(s, "TURNING INSIGHT INTO CONVERSION", "Business Value for Sales")
plays = [
    ("Lead with the bundle", "Auto-only converts at 52% vs. 80% bundled. Proactively "
     "offer a bundle on every auto-only quote.", TEAL),
    ("Smart follow-up", "Model recall of 91% flags warm prospects. Prioritize agent "
     "outreach on high-probability quotes; light-touch the rest.", AMBER),
    ("Set realistic targets", "Segment-level acceptance rates give Sales honest "
     "conversion benchmarks instead of one blanket number.", RUST),
    ("Pricing follow-on", "Add the premium/price field to enable a price-range model "
     "for margin optimization — the natural next dataset.", NAVY),
]
x = Inches(0.6)
for t, d, col in plays:
    card(s, x, Inches(1.6), Inches(2.92), Inches(2.9), t, [d], col)
    x += Inches(3.04)
card(s, Inches(0.6), Inches(4.75), Inches(12.1), Inches(1.7),
     "Estimated impact",
     ["If bundle prompts and prioritized follow-up convert even 3–5% of today's "
      "auto-only decliners, that flows straight to written premium — at near-zero "
      "incremental cost, using rules the team already understands. "
      "(Estimated — to be validated with an A/B test.)"], TEAL)
footer(s, 11)

# ---------------------------------------------------------- 12. Other use cases
s = prs.slides.add_slide(BLANK)
header(s, "WHERE ELSE THIS PAYS OFF", "Similar Use Cases Across the Company")
rows = [
    ("Renewal / churn prediction", "Flag policies likely to lapse and intervene early",
     "Value: retained premium  ·  Obstacle: needs behavioral + billing history"),
    ("Claims fraud triage", "Rank claims by anomaly for investigator focus",
     "Value: loss reduction  ·  Obstacle: labeled fraud cases are scarce"),
    ("Dynamic price-to-convert", "Recommend a quote price that balances win-rate & margin",
     "Value: higher margin  ·  Obstacle: requires price & competitor data"),
    ("Cross-sell propensity", "Predict who'll add home/life to an auto policy",
     "Value: revenue per customer  ·  Obstacle: data silos across product lines"),
]
y = Inches(1.65)
for t, d, note in rows:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(1.18), LIGHT)
    rect(s, Inches(0.6), y, Inches(0.1), Inches(1.18), TEAL)
    textbox(s, Inches(0.9), y + Inches(0.12), Inches(11.6), Inches(1.0),
            [{"text": t, "size": 16, "bold": True, "color": NAVY, "space_after": 1},
             {"text": d, "size": 13, "color": GREY, "space_after": 1},
             {"text": note, "size": 12, "color": RUST}])
    y += Inches(1.3)
footer(s, 12)

# ---------------------------------------------------------- 13. Takeaways / next steps
s = prs.slides.add_slide(BLANK)
header(s, "WHAT TO TELL STAKEHOLDERS", "Key Takeaways & Next Steps")
card(s, Inches(0.6), Inches(1.6), Inches(5.95), Inches(4.0),
     "What we now know", [
        "Conversion is driven mostly by bundling and applicant profile",
        "We can predict acceptance at ~79% accuracy on unseen quotes",
        "A transparent 2-rule model matches the ML — start simple",
        "Error balance favors catching warm leads (91% recall)",
        "No pricing field yet → margin work needs new data",
        "Data is clean, complete, and ready for the next iteration"], TEAL)
card(s, Inches(6.7), Inches(1.6), Inches(6.0), Inches(4.0),
     "Recommended next steps", [
        "Ship bundle-prompt + follow-up rules to Sales now",
        "Run an A/B test to measure real conversion lift",
        "Acquire premium/price data to build a pricing model",
        "Enrich data (credit, tenure, web behavior) before scaling ML",
        "Revisit ML once relationships are genuinely non-linear",
        "Monitor segment acceptance rates as a live dashboard"], NAVY)
footer(s, 13)

# ---------------------------------------------------------- 14. Closing
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.3), SW, Inches(0.08), TEAL)
textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0),
        [{"text": "Start simple. Prove lift. Then scale.", "size": 34,
          "color": WHITE, "bold": True}])
textbox(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.2),
        [{"text": "Two factors explain most of the conversion story — and a simple, "
          "explainable model already captures it. The win is acting on it.",
          "size": 18, "color": RGBColor(0xCF, 0xDD, 0xE8)}])
textbox(s, Inches(0.9), Inches(5.4), Inches(11), Inches(0.8),
        [{"text": "Phillip Smith  ·  Senior ML Engineering  ·  Quote Decision Predictor",
          "size": 14, "color": GREY}])

prs.save(str(OUT))
print("Saved deck:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
